from pptx import Presentation
from pptx.util import Inches
from pathlib import Path


path_cm_folder = Path('/home/local/VANDERBILT/gaoc11/Projects/Variance-Aging-Diffusion/Figure/clustermap_v2')
path_cm_beta = path_cm_folder / 'beta'

prs = Presentation()

list_img_cm_beta = [fn for fn in path_cm_beta.iterdir()]

for img_beta in list_img_cm_beta:
    img_pvalue = path_cm_folder / 'pvalue' / img_beta.name.replace('_beta_','_pvalue_')
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(0)
    top = Inches(0)
    pic = slide.shapes.add_picture(str(img_beta), left, top, width=Inches(10))
    
    top = Inches(3)
    pic_2 = slide.shapes.add_picture(str(img_pvalue), left, top, width=Inches(10))
    
prs.save(path_cm_folder/'clustermap_20230209.pptx')